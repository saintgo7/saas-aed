#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Deck spec(JSON) -> PPTX 렌더러.
사람이 직접 디자인한 듯한 일관된 강의 슬라이드를 생성한다.
사용: python3 build_pptx.py deck.json out.pptx
"""
import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 디자인 토큰 ----
KFONT = "맑은 고딕"
NAVY = RGBColor(0x14, 0x2A, 0x4F)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
ACCENT2 = RGBColor(0xC8, 0x5A, 0x1B)
LIGHT = RGBColor(0xEE, 0xF2, 0xF8)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
LGRAY = RGBColor(0x8A, 0x8A, 0x8A)
DARK = RGBColor(0x23, 0x29, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


def _set_font(run, size=None, bold=None, color=None, italic=None, font=KFONT):
    run.font.name = font
    # 동아시아 폰트도 동일 지정
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _rect(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          wrap=True, space_after=6, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    # runs: list of paragraphs; each paragraph = list of (text, style-dict)
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
        if isinstance(para, tuple):
            para = [para]
        for seg in para:
            txt, st = seg
            r = p.add_run()
            r.text = txt
            _set_font(r, **st)
    return tb


def _footer(slide, idx, course):
    _rect(slide, 0, Inches(7.06), EMU_W, Pt(2), fill=LIGHT)
    _text(slide, Inches(0.55), Inches(7.08), Inches(9), Inches(0.32),
          [[(course, dict(size=9, color=LGRAY))]], anchor=MSO_ANCHOR.MIDDLE)
    _text(slide, Inches(11.6), Inches(7.08), Inches(1.2), Inches(0.32),
          [[(str(idx), dict(size=10, color=GRAY, bold=True))]],
          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _content_header(slide, title, kicker=None):
    _rect(slide, 0, 0, Inches(0.22), EMU_H, fill=NAVY)
    if kicker:
        _text(slide, Inches(0.55), Inches(0.34), Inches(11), Inches(0.3),
              [[(kicker, dict(size=11, color=ACCENT, bold=True))]])
    _text(slide, Inches(0.52), Inches(0.6), Inches(12.2), Inches(0.95),
          [[(title, dict(size=27, color=NAVY, bold=True))]], line_spacing=1.0)
    _rect(slide, Inches(0.57), Inches(1.5), Inches(1.1), Pt(3), fill=ACCENT2)


# ---- 레이아웃들 ----
def slide_title(prs, s, course):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, EMU_W, EMU_H, fill=NAVY)
    _rect(sl, 0, Inches(4.62), EMU_W, Pt(3), fill=ACCENT2)
    _rect(sl, Inches(0.9), Inches(1.05), Inches(1.6), Pt(4), fill=ACCENT)
    if s.get("kicker"):
        _text(sl, Inches(0.9), Inches(1.25), Inches(11), Inches(0.4),
              [[(s["kicker"], dict(size=15, color=RGBColor(0x9D, 0xBE, 0xE6), bold=True))]])
    _text(sl, Inches(0.9), Inches(1.9), Inches(11.5), Inches(2.3),
          [[(s.get("title", ""), dict(size=40, color=WHITE, bold=True))]], line_spacing=1.08)
    if s.get("subtitle"):
        _text(sl, Inches(0.92), Inches(4.85), Inches(11.5), Inches(1.0),
              [[(s["subtitle"], dict(size=18, color=RGBColor(0xCF, 0xDD, 0xEF)))]])
    meta = []
    if s.get("org"):
        meta.append([(s["org"], dict(size=14, color=RGBColor(0xB9, 0xCB, 0xE2)))])
    if s.get("date") or s.get("presenter"):
        line = "   ·   ".join([x for x in [s.get("date", ""), s.get("presenter", "")] if x])
        meta.append([(line, dict(size=13, color=RGBColor(0x8F, 0xA8, 0xC9)))])
    if meta:
        _text(sl, Inches(0.92), Inches(6.2), Inches(11.5), Inches(0.9), meta, space_after=3)


def slide_section(prs, s, course, num):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, EMU_W, EMU_H, fill=LIGHT)
    _rect(sl, 0, 0, Inches(4.4), EMU_H, fill=NAVY)
    _text(sl, Inches(0.7), Inches(2.5), Inches(3.2), Inches(2.0),
          [[(s.get("num", num), dict(size=80, color=RGBColor(0x3E, 0x5C, 0x8A), bold=True))]])
    _rect(sl, Inches(4.9), Inches(3.0), Inches(0.7), Pt(4), fill=ACCENT2)
    _text(sl, Inches(4.9), Inches(3.25), Inches(7.8), Inches(1.8),
          [[(s.get("title", ""), dict(size=32, color=NAVY, bold=True))]], line_spacing=1.05)
    if s.get("subtitle"):
        _text(sl, Inches(4.92), Inches(4.7), Inches(7.6), Inches(1.2),
              [[(s["subtitle"], dict(size=15, color=GRAY))]])


def _bullets_tf(slide, x, y, w, h, bullets, base=16):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    first = True
    for b in bullets:
        if isinstance(b, str):
            b = {"text": b, "level": 0}
        lvl = b.get("level", 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = min(lvl, 2)
        p.space_after = Pt(7 if lvl == 0 else 4)
        p.line_spacing = 1.08
        mark = "▪  " if lvl == 0 else ("–  " if lvl == 1 else "·  ")
        r = p.add_run(); r.text = mark
        _set_font(r, size=base - lvl * 2,
                  color=(ACCENT2 if lvl == 0 else LGRAY), bold=(lvl == 0))
        r2 = p.add_run(); r2.text = b["text"]
        _set_font(r2, size=base - lvl * 2,
                  color=(DARK if lvl == 0 else GRAY),
                  bold=b.get("bold", False))
    return tb


def slide_bullets(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    body_y = Inches(1.85)
    if s.get("lead"):
        _text(sl, Inches(0.6), body_y, Inches(12.1), Inches(0.7),
              [[(s["lead"], dict(size=15, color=ACCENT, italic=True))]])
        body_y = Inches(2.5)
    _bullets_tf(sl, Inches(0.62), body_y, Inches(12.1), Inches(4.4),
                s.get("bullets", []), base=s.get("font", 17))
    _footer(sl, idx, course)


def slide_two_col(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    cols = [("left", Inches(0.62)), ("right", Inches(6.95))]
    for key, x in cols:
        col = s.get(key, {})
        head = col.get("title", "")
        _rect(sl, x, Inches(1.95), Inches(5.75), Inches(0.55), fill=NAVY)
        _text(sl, x + Inches(0.18), Inches(1.97), Inches(5.4), Inches(0.5),
              [[(head, dict(size=15, color=WHITE, bold=True))]], anchor=MSO_ANCHOR.MIDDLE)
        _rect(sl, x, Inches(2.5), Inches(5.75), Inches(4.2), fill=LIGHT)
        _bullets_tf(sl, x + Inches(0.2), Inches(2.7), Inches(5.4), Inches(3.9),
                    col.get("bullets", []), base=15)
    _footer(sl, idx, course)


def slide_table(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    data = s.get("table", [])
    if not data:
        _footer(sl, idx, course); return
    rows, cols = len(data), len(data[0])
    top = Inches(2.0)
    h = min(Inches(4.6), Inches(0.55) * rows)
    gtbl = sl.shapes.add_table(rows, cols, Inches(0.62), top, Inches(12.1), h).table
    for ci in range(cols):
        gtbl.columns[ci].width = Emu(int(Inches(12.1) / cols))
    for ri, row in enumerate(data):
        gtbl.rows[ri].height = Emu(int(h / rows))
        for ci in range(cols):
            cell = gtbl.cell(ri, ci)
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if ri == 0 else (WHITE if ri % 2 else LIGHT)
            txt = str(row[ci]) if ci < len(row) else ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = txt
            _set_font(r, size=12 if ri else 13,
                      color=(WHITE if ri == 0 else DARK), bold=(ri == 0))
    _footer(sl, idx, course)


def slide_takeaway(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, EMU_W, EMU_H, fill=NAVY)
    _rect(sl, Inches(1.0), Inches(1.4), Inches(11.33), Inches(4.7), fill=RGBColor(0x1C, 0x37, 0x63))
    _rect(sl, Inches(1.0), Inches(1.4), Inches(0.16), Inches(4.7), fill=ACCENT2)
    _text(sl, Inches(1.5), Inches(1.85), Inches(10.4), Inches(0.6),
          [[(s.get("label", "핵심 정리"), dict(size=15, color=ACCENT, bold=True))]])
    _text(sl, Inches(1.5), Inches(2.55), Inches(10.4), Inches(3.2),
          [[(s.get("text", ""), dict(size=27, color=WHITE, bold=True))]], line_spacing=1.18)
    _footer(sl, idx, course)


def slide_quote(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, EMU_W, EMU_H, fill=LIGHT)
    _text(sl, Inches(1.1), Inches(1.5), Inches(2), Inches(1.2),
          [[("“", dict(size=90, color=ACCENT, bold=True))]])
    _text(sl, Inches(1.4), Inches(2.7), Inches(10.5), Inches(2.8),
          [[(s.get("text", ""), dict(size=26, color=NAVY, bold=True))]], line_spacing=1.2)
    if s.get("by"):
        _text(sl, Inches(1.45), Inches(5.6), Inches(10), Inches(0.5),
              [[("— " + s["by"], dict(size=15, color=GRAY))]])
    _footer(sl, idx, course)


def slide_closing(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(sl, 0, 0, EMU_W, EMU_H, fill=NAVY)
    _rect(sl, Inches(0.9), Inches(2.4), Inches(1.6), Pt(4), fill=ACCENT2)
    _text(sl, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.4),
          [[(s.get("title", "감사합니다"), dict(size=40, color=WHITE, bold=True))]])
    lines = [[(l, dict(size=16, color=RGBColor(0xC5, 0xD5, 0xEA)))] for l in s.get("lines", [])]
    if lines:
        _text(sl, Inches(0.92), Inches(4.3), Inches(11.5), Inches(2.2), lines, space_after=5)


# =====================================================================
#  그림(다이어그램) 레이아웃 — 이미지가 아닌 편집 가능한 네이티브 도형
# =====================================================================
DIA_FILLS = [RGBColor(0x21, 0x3B, 0x66), RGBColor(0x2E, 0x75, 0xB6),
             RGBColor(0x57, 0x9A, 0xD6), RGBColor(0x9D, 0xC3, 0xE6)]


def _box(slide, x, y, w, h, title=None, desc=None, fill=LIGHT, edge=None,
         tcolor=NAVY, dcolor=GRAY, tsize=14, dsize=11, accent=None,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, radius=True):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = _rect(slide, x, y, w, h, fill=fill, line=edge, shape=shape)
    if accent:
        _rect(slide, x, y, w, Pt(6), fill=accent)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = align
    if title:
        r = p.add_run(); r.text = title
        _set_font(r, size=tsize, bold=True, color=tcolor)
    if desc:
        p2 = tf.add_paragraph(); p2.alignment = align; p2.space_before = Pt(2)
        p2.line_spacing = 1.02
        r2 = p2.add_run(); r2.text = desc
        _set_font(r2, size=dsize, color=dcolor)
    return sp


def _arrow(slide, x, y, w, h, d="r", color=ACCENT2):
    shp = {"r": MSO_SHAPE.RIGHT_ARROW, "l": MSO_SHAPE.LEFT_ARROW,
           "u": MSO_SHAPE.UP_ARROW, "d": MSO_SHAPE.DOWN_ARROW}[d]
    return _rect(slide, x, y, w, h, fill=color, shape=shp)


def _badge(slide, x, y, d, text, fill=ACCENT2, tcolor=WHITE, size=18):
    sp = _rect(slide, x, y, d, d, fill=fill, shape=MSO_SHAPE.OVAL)
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    _set_font(r, size=size, bold=True, color=tcolor)
    return sp


def slide_flow(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    steps = s.get("steps", [])
    n = max(1, len(steps))
    vertical = s.get("dir") == "v" or n > 5
    if not vertical:
        aw = Inches(0.55)
        x0 = Inches(0.62)
        total = Inches(12.1)
        bw = int((total - (n - 1) * aw) / n)
        by = Inches(3.0); bh = Inches(2.0)
        for i, st in enumerate(steps):
            bx = x0 + i * (bw + aw)
            _box(sl, bx, by, bw, bh, st.get("t"), st.get("d"),
                 fill=LIGHT, accent=DIA_FILLS[i % len(DIA_FILLS)],
                 tcolor=NAVY, tsize=15, dsize=11)
            if i < n - 1:
                _arrow(sl, bx + bw + Inches(0.05), by + int(bh / 2) - Inches(0.22),
                       Inches(0.45), Inches(0.44), "r")
    else:
        x0 = Inches(3.0); bw = Inches(7.3)
        ah = Inches(0.4)
        top = Inches(1.95)
        bh = int((Inches(4.95) - (n - 1) * ah) / n)
        bh = min(bh, Inches(1.1))
        for i, st in enumerate(steps):
            by = top + i * (bh + ah)
            _box(sl, x0, by, bw, bh, st.get("t"), st.get("d"),
                 fill=LIGHT, accent=DIA_FILLS[i % len(DIA_FILLS)],
                 align=PP_ALIGN.LEFT, tsize=14, dsize=11)
            if i < n - 1:
                _arrow(sl, x0 + int(bw / 2) - Inches(0.2), by + bh + Inches(0.02),
                       Inches(0.4), Inches(0.36), "d")
    _footer(sl, idx, course)


def slide_loop(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    nodes = s.get("nodes", [])[:4]
    bw = Inches(3.7); bh = Inches(1.35)
    L, R = Inches(0.9), Inches(8.72)
    T, B = Inches(2.35), Inches(5.1)
    pos = [(L, T), (R, T), (R, B), (L, B)]
    for i, nd in enumerate(nodes):
        _box(sl, pos[i][0], pos[i][1], bw, bh, nd.get("t"), nd.get("d"),
             fill=DIA_FILLS[i % len(DIA_FILLS)], tcolor=WHITE, dcolor=RGBColor(0xE5, 0xEE, 0xF8),
             tsize=15, dsize=11)
    cx, cy = Inches(5.9), Inches(3.55)
    _badge(sl, cx, cy, Inches(1.5), s.get("center", "반복"),
           fill=ACCENT2, size=15)
    midy = T + int(bh / 2) - Inches(0.2)
    _arrow(sl, L + bw + Inches(0.2), midy, Inches(1.0), Inches(0.42), "r")
    _arrow(sl, R + int(bw / 2) - Inches(0.2), T + bh + Inches(0.12), Inches(0.42), Inches(0.95), "d")
    _arrow(sl, L + bw + Inches(0.2), B + int(bh / 2) - Inches(0.2), Inches(1.0), Inches(0.42), "l")
    _arrow(sl, L + int(bw / 2) - Inches(0.2), T + bh + Inches(0.12), Inches(0.42), Inches(0.95), "u")
    _footer(sl, idx, course)


def slide_layers(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    items = s.get("items", [])[:4]
    n = len(items)
    ox, oy = Inches(2.1), Inches(2.15)
    ow, oh = Inches(9.1), Inches(4.55)
    dx, dy = Inches(0.95), Inches(0.62)
    for i, it in enumerate(items):
        x = ox + i * dx; y = oy + i * dy
        w = ow - 2 * i * dx; h = oh - 2 * i * dy
        _rect(sl, x, y, w, h, fill=DIA_FILLS[i % len(DIA_FILLS)],
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb = _text(sl, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.45),
                   [[(it.get("t", ""), dict(size=15 - i, bold=True,
                      color=(WHITE if i < 2 else NAVY)))]], align=PP_ALIGN.CENTER,
                   anchor=MSO_ANCHOR.TOP)
    # 우측 설명 캡션
    cap = [[(f"{it.get('t','')} — {it.get('d','')}",
             dict(size=11, color=GRAY))] for it in items]
    if any(it.get("d") for it in items):
        _text(sl, Inches(0.62), Inches(6.55), Inches(12.1), Inches(0.4),
              [[("바깥에서 안으로: " + "  ›  ".join(it.get("t", "") for it in items),
                 dict(size=11, color=LGRAY, italic=True))]])
    _footer(sl, idx, course)


def slide_blocks(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    items = s.get("items", [])
    cols = s.get("cols", 2 if len(items) <= 4 else 3)
    rows = (len(items) + cols - 1) // cols
    gx, gy = Inches(0.35), Inches(0.35)
    x0, y0 = Inches(0.62), Inches(2.0)
    aw, ah = Inches(12.1), Inches(4.7)
    cw = int((aw - (cols - 1) * gx) / cols)
    ch = int((ah - (rows - 1) * gy) / rows)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        _box(sl, x, y, cw, ch, None, None, fill=LIGHT,
             accent=DIA_FILLS[i % len(DIA_FILLS)], radius=True)
        bx = x + Inches(0.2); by = y + Inches(0.22)
        if it.get("n"):
            _badge(sl, bx, by, Inches(0.5), str(it["n"]),
                   fill=DIA_FILLS[i % len(DIA_FILLS)], size=15)
            tx = bx + Inches(0.68)
        else:
            tx = bx
        _text(sl, tx, by + Inches(0.02), x + cw - tx - Inches(0.15), Inches(0.5),
              [[(it.get("t", ""), dict(size=15, bold=True, color=NAVY))]])
        if it.get("d"):
            _text(sl, bx, by + Inches(0.62), cw - Inches(0.4), ch - Inches(0.9),
                  [[(it["d"], dict(size=12, color=GRAY))]], line_spacing=1.05)
    _footer(sl, idx, course)


def slide_steps(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    steps = s.get("steps", [])
    n = max(1, len(steps))
    top = Inches(2.0)
    rh = min(Inches(0.92), int((Inches(4.9)) / n))
    gap = Inches(0.12)
    for i, st in enumerate(steps):
        y = top + i * (rh + gap)
        _badge(sl, Inches(0.7), y + int(rh / 2) - Inches(0.3), Inches(0.6),
               str(i + 1), fill=DIA_FILLS[i % len(DIA_FILLS)], size=18)
        _rect(sl, Inches(1.5), y, Inches(11.2), rh, fill=LIGHT,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(sl, Inches(1.75), y + Inches(0.08), Inches(10.7), rh - Inches(0.12),
              [[(st.get("t", ""), dict(size=14, bold=True, color=NAVY)),
                ("   " + st.get("d", ""), dict(size=12, color=GRAY))]],
              anchor=MSO_ANCHOR.MIDDLE)
    _footer(sl, idx, course)


def slide_compare(prs, s, course, idx):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _content_header(sl, s.get("title", ""), s.get("kicker"))
    panels = [("left", Inches(0.62), DIA_FILLS[0]), ("right", Inches(6.95), DIA_FILLS[1])]
    for key, x, col in panels:
        c = s.get(key, {})
        _rect(sl, x, Inches(2.0), Inches(5.76), Inches(0.62), fill=col,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _text(sl, x + Inches(0.2), Inches(2.02), Inches(5.4), Inches(0.58),
              [[(c.get("title", ""), dict(size=16, bold=True, color=WHITE))]],
              anchor=MSO_ANCHOR.MIDDLE)
        _rect(sl, x, Inches(2.72), Inches(5.76), Inches(3.95), fill=LIGHT,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _bullets_tf(sl, x + Inches(0.25), Inches(2.95), Inches(5.3), Inches(3.5),
                    c.get("bullets", []), base=14)
    _badge(sl, Inches(6.16), Inches(3.95), Inches(1.0), s.get("vs", "VS"),
           fill=ACCENT2, size=18)
    _footer(sl, idx, course)


LAYOUTS = {
    "title": slide_title, "section": slide_section, "bullets": slide_bullets,
    "two_col": slide_two_col, "table": slide_table, "takeaway": slide_takeaway,
    "quote": slide_quote, "closing": slide_closing,
    "flow": slide_flow, "loop": slide_loop, "layers": slide_layers,
    "blocks": slide_blocks, "steps": slide_steps, "compare": slide_compare,
}


def build(deck, out):
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    course = deck.get("course", "")
    sec_num = 0
    for i, s in enumerate(deck["slides"], start=1):
        lay = s.get("layout", "bullets")
        fn = LAYOUTS.get(lay, slide_bullets)
        if lay == "section":
            sec_num += 1
            sl_before = len(prs.slides._sldIdLst)
            fn(prs, s, course, str(sec_num).zfill(2))
        elif lay in ("title",):
            fn(prs, s, course)
        else:
            fn(prs, s, course, i)
        # 발표자 노트
        notes = s.get("notes")
        if notes:
            slide = prs.slides[-1]
            slide.notes_slide.notes_text_frame.text = notes
    # 문서 속성(사람 작성처럼)
    cp = prs.core_properties
    cp.author = deck.get("author", "특장차인증센터 교육운영팀")
    cp.last_modified_by = deck.get("author", "특장차인증센터 교육운영팀")
    cp.title = deck.get("title", course)
    cp.subject = deck.get("subject", "AI 활용 실무 특강")
    cp.category = "교육자료"
    cp.comments = ""
    prs.save(out)
    print(f"OK {out} ({len(deck['slides'])} slides)")


if __name__ == "__main__":
    with open(sys.argv[1], encoding="utf-8") as f:
        deck = json.load(f)
    build(deck, sys.argv[2])
